"""ImpactAI Foundry — redesign design system.
Warm-minimal *indic* palette · bold modern sans · strict grid.
All geometry in inches; one shared grid so every slide aligns.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# CANVAS  (16:9 widescreen)
# ----------------------------------------------------------------------------
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
PAGE_W = 13.333
PAGE_H = 7.5

# Grid — one set of margins for the entire deck
MARGIN_L = 0.92
MARGIN_R = 0.92
MARGIN_T = 0.72
MARGIN_B = 0.62
CONTENT_W = PAGE_W - MARGIN_L - MARGIN_R   # 11.493
CONTENT_R = PAGE_W - MARGIN_R

# ----------------------------------------------------------------------------
# PALETTE  (warm-minimal indic)
# ----------------------------------------------------------------------------
INDIGO      = RGBColor(0x1F, 0x2A, 0x52)   # deep indian indigo — primary dark
INDIGO_DEEP = RGBColor(0x16, 0x1F, 0x3D)   # full-bleed dark fields
SAFFRON     = RGBColor(0xF2, 0xA5, 0x16)   # marigold — primary accent
SAFFRON_DK  = RGBColor(0xDE, 0x8A, 0x12)
TERRACOTTA  = RGBColor(0xBE, 0x4A, 0x2E)   # secondary accent (sparse)
SPICE       = RGBColor(0xD2, 0x69, 0x1E)   # burnt orange (sequence)
SAND        = RGBColor(0xB7, 0x9A, 0x57)   # muted gold
TEAL        = RGBColor(0x2E, 0x7A, 0x6B)   # rare cool support

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAM  = RGBColor(0xFB, 0xF6, 0xEC)        # warm off-white bg
PAPER  = RGBColor(0xF4, 0xEE, 0xDF)
INK    = RGBColor(0x1C, 0x1C, 0x24)        # body text on light
MUTE   = RGBColor(0x6F, 0x6A, 0x60)        # secondary text / eyebrow
LINE   = RGBColor(0xE6, 0xE0, 0xD2)        # hairline on light
LINE_W = RGBColor(0x2E, 0x37, 0x5E)        # hairline on indigo (lighter)
CREAM_DIM = RGBColor(0x9A, 0xA0, 0xB8)     # muted text on indigo

# ----------------------------------------------------------------------------
# TYPE
# ----------------------------------------------------------------------------
DISPLAY = "Poppins"            # SemiBold/Bold for titles
DISPLAY_SB = "Poppins SemiBold"
BODY    = "Inter"              # body
EYEBROW_FONT = "Inter"


# ----------------------------------------------------------------------------
# LOW-LEVEL HELPERS
# ----------------------------------------------------------------------------
def _no_line(shape):
    shape.line.fill.background()


def _no_shadow(shape):
    # add_shape attaches <p:style> with an effectRef -> theme drop-shadow.
    # Remove it so shapes are purely flat (we always set fill/line explicitly).
    el = shape._element
    for child in list(el):
        if child.tag.endswith('}style'):
            el.remove(child)
    shape.shadow.inherit = False


def bg(slide, color):
    """Full-bleed background rectangle (added first => sits at back)."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    _no_line(r); _no_shadow(r)
    return r


def rect(slide, x, y, w, h, color, line=None, line_w=0.75):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        r.fill.background()
    else:
        r.fill.solid(); r.fill.fore_color.rgb = color
    if line is None:
        _no_line(r)
    else:
        r.line.color.rgb = line; r.line.width = Pt(line_w)
    _no_shadow(r)
    return r


def rule(slide, x, y, w, color, weight=3.0):
    """A short accent rule (the signature saffron mark)."""
    return rect(slide, x, y, w, weight / 72.0, color)


def set_spc(run, pts):
    """Letter-spacing (tracking) in points."""
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(round(pts * 100))))


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def para(tf, first=False, align=PP_ALIGN.LEFT, before=0, after=0, line=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = line
    return p


def run(p, text, font=BODY, size=18, color=INK, bold=False, italic=False, spc=None):
    r = p.add_run(); r.text = text
    f = r.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    if spc is not None:
        set_spc(r, spc)
    return r


# ----------------------------------------------------------------------------
# COMPONENTS  (the through-line that keeps every slide aligned)
# ----------------------------------------------------------------------------
def eyebrow(slide, text, color=SAFFRON_DK, x=MARGIN_L, y=MARGIN_T, tick=True):
    """All-caps tracked eyebrow with optional saffron tick mark."""
    tx = x
    if tick:
        rect(slide, x, y + 0.045, 0.30, 0.075, SAFFRON)
        tx = x + 0.44
    tb, tf = textbox(slide, tx, y - 0.04, CONTENT_W - (tx - x), 0.35)
    p = para(tf, first=True)
    run(p, text.upper(), font=EYEBROW_FONT, size=12.5, color=color, bold=True, spc=2.2)
    return tb


def title(slide, text, x=MARGIN_L, y=1.28, w=CONTENT_W, size=40, color=INDIGO,
          underline=True, underline_w=0.62):
    tb, tf = textbox(slide, x, y, w, 2.2)
    p = para(tf, first=True, line=1.02)
    run(p, text, font=DISPLAY, size=size, color=color, bold=True)
    if underline:
        # placed by caller via title_rule when exact y is known
        pass
    return tb


def title_rule(slide, x=MARGIN_L, y=None, w=0.62, color=SAFFRON):
    return rule(slide, x, y, w, color, weight=3.5)


def footer(slide, page=None, dark=False):
    col = CREAM_DIM if dark else MUTE
    tb, tf = textbox(slide, MARGIN_L, PAGE_H - 0.52, 6.0, 0.3)
    p = para(tf, first=True)
    run(p, "IMPACTAI FOUNDRY", font=EYEBROW_FONT, size=8.5, color=col, bold=True, spc=1.6)
    run(p, "   ·   BANGALORE 2026", font=EYEBROW_FONT, size=8.5, color=col, bold=False, spc=1.6)
    if page is not None:
        tb2, tf2 = textbox(slide, CONTENT_R - 1.0, PAGE_H - 0.52, 1.0, 0.3)
        p2 = para(tf2, first=True, align=PP_ALIGN.RIGHT)
        run(p2, str(page), font=EYEBROW_FONT, size=8.5, color=col, bold=True, spc=1.0)


def corner_motif(slide, color=SAFFRON, x=None, y=None, opacity_layer=False):
    """Minimal 'indic' identity: nested right-angle ticks (jaali corner)."""
    if x is None:
        x = CONTENT_R - 0.9
    if y is None:
        y = MARGIN_T - 0.06
    for i in range(3):
        off = i * 0.16
        # horizontal stub
        rect(slide, x + off, y, 0.5 - off, 0.028, color)
        # vertical stub
        rect(slide, x + 0.5 - 0.028, y, 0.028, 0.5 - off, color)
