"""ImpactAI Foundry — INDIC design system.
Cream parchment · razor-sharp kilim corner motifs · Georgia serif ·
navy / terracotta / gold / green. Decoded-Futures-clean layout discipline.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- canvas (16:9) ----
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
PAGE_W, PAGE_H = 13.333, 7.5

# ---- grid: content sits INSIDE the corner frame ----
MARGIN_L = 1.55
MARGIN_R = 1.55
MARGIN_T = 1.18
MARGIN_B = 1.0
CONTENT_W = PAGE_W - MARGIN_L - MARGIN_R
CONTENT_R = PAGE_W - MARGIN_R

# ---- palette ----
NAVY      = RGBColor(0x1F, 0x3A, 0x5F)
NAVY_DEEP = RGBColor(0x17, 0x2C, 0x49)   # dark slide bg
INK       = RGBColor(0x33, 0x31, 0x2B)   # warm near-black body
TERRA     = RGBColor(0xB8, 0x50, 0x42)
TERRA_DK  = RGBColor(0xA0, 0x3D, 0x30)
GOLD      = RGBColor(0xC8, 0xA0, 0x4A)
GOLD_DK   = RGBColor(0xA8, 0x84, 0x2F)
GOLDM     = RGBColor(0x8A, 0x7A, 0x55)
GREEN     = RGBColor(0x4A, 0x6B, 0x3A)
CREAM     = RGBColor(0xF4, 0xEC, 0xD8)
PAPER     = RGBColor(0xEC, 0xE2, 0xC8)   # card fill
LINE      = RGBColor(0xD8, 0xCB, 0xA8)   # hairline on cream
CREAM_TXT = RGBColor(0xF1, 0xE7, 0xCE)   # text on navy
CREAM_DIM = RGBColor(0xC9, 0xBA, 0x97)   # muted text on navy

SERIF = "Georgia"

# ---- corner art: the ORIGINAL ornate kilim corner, extracted from the
#      source deck's parchment border (assets/kilim_source.png) ----
import extract_corner
_CORNER_PATH = "corner_real.png"
_CORNER_DARK_PATH = "corner_real_dark.png"
CORNER_ASPECT = 289 / 263.0   # w/h of the extracted ornament


def ensure_assets():
    extract_corner.make_corners()


# ---------------------------------------------------------------- primitives
def _flat(shape):
    el = shape._element
    for ch in list(el):
        if ch.tag.endswith('}style'):
            el.remove(ch)
    shape.shadow.inherit = False


def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); _flat(r)
    return r


def rect(slide, x, y, w, h, color, line=None, line_w=0.75):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        r.fill.background()
    else:
        r.fill.solid(); r.fill.fore_color.rgb = color
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(line_w)
    _flat(r)
    return r


def rule(slide, x, y, w, color, weight=2.5):
    return rect(slide, x, y, w, weight / 72.0, color)


def set_spc(run, pts):
    run._r.get_or_add_rPr().set('spc', str(int(round(pts * 100))))


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def para(tf, first=False, align=PP_ALIGN.LEFT, before=0, after=0, line=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = line
    return p


def run(p, text, font=SERIF, size=18, color=INK, bold=False, italic=False, spc=None):
    r = p.add_run(); r.text = text
    f = r.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    if spc is not None:
        set_spc(r, spc)
    return r


# ---------------------------------------------------------------- corners
def corners(slide, size=1.45, off=0.42, which=("tl", "br"), dark=False):
    """Place the ornate kilim corner ornament. `size` = WIDTH in inches;
    height follows the art's true aspect. Default: top-left + bottom-right
    diagonal pair (the user's 'one left, one right opposite corners')."""
    path = _CORNER_DARK_PATH if dark else _CORNER_PATH
    w = size
    h = size / CORNER_ASPECT
    spots = {
        "tl": (off, off, False, False),
        "tr": (PAGE_W - off - w, off, True, False),
        "bl": (off, PAGE_H - off - h, False, True),
        "br": (PAGE_W - off - w, PAGE_H - off - h, True, True),
    }
    for k in which:
        x, y, fh, fv = spots[k]
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        xfrm = pic._element.find('.//' + qn('a:xfrm'))
        if fh:
            xfrm.set('flipH', '1')
        if fv:
            xfrm.set('flipV', '1')
    return slide


# ---------------------------------------------------------------- components
def eyebrow(slide, text, x=MARGIN_L, y=MARGIN_T, color=GOLD_DK, align=PP_ALIGN.LEFT, w=None):
    if w is None:
        w = CONTENT_W
    tb, tf = textbox(slide, x, y, w, 0.34)
    p = para(tf, first=True, align=align)
    run(p, text.upper(), font=SERIF, size=13, color=color, bold=True, spc=2.6)
    return tb


def footer(slide, page=None, dark=False):
    col = CREAM_DIM if dark else GOLDM
    tb, tf = textbox(slide, MARGIN_L, PAGE_H - 0.62, 7.0, 0.3)
    p = para(tf, first=True)
    run(p, "IMPACTAI FOUNDRY", font=SERIF, size=8.5, color=col, bold=True, spc=2.0)
    run(p, "   ·   BANGALORE 2026", font=SERIF, size=8.5, color=col, spc=2.0)
    if page is not None:
        tb2, tf2 = textbox(slide, CONTENT_R - 1.0, PAGE_H - 0.62, 1.0, 0.3)
        p2 = para(tf2, first=True, align=PP_ALIGN.RIGHT)
        run(p2, str(page), font=SERIF, size=8.5, color=col, bold=True)
