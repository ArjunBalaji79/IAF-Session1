"""Archetype renderer library for the IAF Session-1 redesign.

Each renderer takes (slide, rec) where rec is a plain dict describing the
slide's content. The build script (deck_build.py) maps the original deck's
content into these records and calls render(slide, rec).

Design system + primitives live in theme_indic. Corners use the REAL
extracted kilim motif (corner_real.png), tl+br by default, tl+tr for
bottom-heavy layouts (see corners_for()).
"""
import os
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
from theme_indic import (
    EMU_W, EMU_H, PAGE_W, PAGE_H, MARGIN_L, MARGIN_R, MARGIN_T, MARGIN_B,
    CONTENT_W, CONTENT_R, NAVY, NAVY_DEEP, INK, TERRA, TERRA_DK, GOLD,
    GOLD_DK, GOLDM, GREEN, CREAM, PAPER, LINE, CREAM_TXT, CREAM_DIM, SERIF,
    bg, rect, rule, textbox, para, run, eyebrow as _eyebrow, footer, corners,
    set_spc,
)

MEDIA = "assets/media"
SPICE = RGBColor(0x9E, 0x6B, 0x33)
GHOST = RGBColor(0xEF, 0xE7, 0xD2)   # faint watermark numeral (lighter than PAPER)

import json as _json
try:
    _REMAP = _json.load(open("media_remap.json"))
except Exception:
    _REMAP = {}

COLORMAP = {
    "navy": NAVY, "terra": TERRA, "gold": GOLD, "green": GREEN,
    "spice": SPICE, "ink": INK, "cream": CREAM_TXT, "gold_dk": GOLD_DK,
}


def _col(name, default=NAVY):
    if isinstance(name, RGBColor):
        return name
    return COLORMAP.get(name, default)


# ---------------------------------------------------------------- shared bits
import math


def est_lines(text, size, width_in, char_factor=0.56):
    """Conservatively estimate wrapped line count for a Georgia run.
    Biased to slightly OVER-count (reserve more vertical space, never collide)."""
    if not text:
        return 1
    char_w = char_factor * size / 72.0
    cpl = max(1, int(width_in / char_w))
    return max(1, math.ceil(len(text) / cpl))


def _title_block(slide, rec, y_eyebrow=MARGIN_T, title_size=33, dark=False):
    """Eyebrow + title + gold rule, left-aligned. Returns y after the rule."""
    tcol = CREAM_TXT if dark else NAVY
    ecol = GOLD if dark else GOLD_DK
    eb = rec.get("eyebrow")
    y = y_eyebrow
    if eb:
        _eyebrow(slide, eb, x=MARGIN_L, y=y, color=ecol)
        y += 0.38
    title = rec.get("title", "")
    lines = est_lines(title, title_size, CONTENT_W)
    tb, tf = textbox(slide, MARGIN_L, y, CONTENT_W, max(1.4, lines * 0.7))
    p = para(tf, first=True, line=1.04)
    run(p, title, font=SERIF, size=title_size, color=tcol, bold=True)
    y += 0.14 + lines * (title_size / 72.0) * 1.26
    rule(slide, MARGIN_L, y, 0.62, GOLD, weight=2.8)
    return y + 0.30


def _bullets(slide, items, x, y, w, h, size=16.5, color=INK, gap=9,
             marker=TERRA, line=1.22, anchor=MSO_ANCHOR.TOP):
    tb, tf = textbox(slide, x, y, w, h, anchor=anchor)
    for i, it in enumerate(items):
        # support {"t":..,"lvl":..} or plain string
        if isinstance(it, dict):
            txt = it.get("t", ""); lvl = it.get("lvl", 0)
        else:
            txt = it; lvl = 0
        p = para(tf, first=(i == 0), line=line, after=gap)
        if lvl and lvl > 0:
            p.level = min(lvl, 4)
            run(p, "–  ", font=SERIF, size=size - 1, color=GOLDM, bold=True)
            run(p, txt, font=SERIF, size=size - 1, color=color)
        else:
            run(p, "•  ", font=SERIF, size=size, color=marker, bold=True)
            run(p, txt, font=SERIF, size=size, color=color)
    return tb


def _img_aspect(path):
    try:
        with Image.open(path) as im:
            return im.size[0] / im.size[1]
    except Exception:
        return 1.6


def place_image(slide, fname, x, y, w, h, fit="contain", frame=True,
                frame_color=None):
    """Place an image fitted into box (x,y,w,h) in inches.
    fit='contain' letterboxes (centered); fit='cover' fills+crops."""
    fname = _REMAP.get(fname, fname)
    path = os.path.join(MEDIA, fname)
    if not os.path.exists(path):
        # placeholder
        rect(slide, x, y, w, h, PAPER, line=LINE, line_w=1.0)
        return
    ar = _img_aspect(path)         # w/h
    box_ar = w / h
    if fit == "cover":
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y),
                                       Inches(w), Inches(h))
        # crop to box aspect
        if ar > box_ar:
            crop = (1 - box_ar / ar) / 2
            pic.crop_left = crop; pic.crop_right = crop
        else:
            crop = (1 - ar / box_ar) / 2
            pic.crop_top = crop; pic.crop_bottom = crop
        iw, ih, ix, iy = w, h, x, y
    else:  # contain
        if ar > box_ar:
            iw = w; ih = w / ar
        else:
            ih = h; iw = h * ar
        ix = x + (w - iw) / 2
        iy = y + (h - ih) / 2
        pic = slide.shapes.add_picture(path, Inches(ix), Inches(iy),
                                       Inches(iw), Inches(ih))
    if frame:
        fc = frame_color or LINE
        ln = rect(slide, ix, iy, iw, ih, None, line=fc, line_w=1.25)
    return pic


# ================================================================ ARCHETYPES
def r_cover(slide, rec):
    bg(slide, NAVY_DEEP if rec.get("dark") else CREAM)
    cx = PAGE_W / 2
    if rec.get("eyebrow"):
        tb, tf = textbox(slide, 0, 2.42, PAGE_W, 0.4)
        p = para(tf, first=True, align=PP_ALIGN.CENTER)
        parts = rec["eyebrow"]
        if isinstance(parts, str):
            run(p, parts, font=SERIF, size=13.5, color=TERRA, bold=True, spc=3.0)
        else:
            for seg in parts:
                run(p, seg.get("t", ""), font=SERIF, size=13.5,
                    color=_col(seg.get("c", "gold"), GOLDM), bold=True, spc=3.0)
    ctitle = rec.get("title", "")
    csize = rec.get("title_size", 60)
    # shrink to keep the cover title on one clean line (floor 40pt)
    while csize > 40 and est_lines(ctitle, csize, PAGE_W - 1.4) > 1:
        csize -= 2
    tb, tf = textbox(slide, 0.6, 2.84, PAGE_W - 1.2, 1.5)
    p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.0)
    run(p, ctitle, font=SERIF, size=csize,
        color=CREAM_TXT if rec.get("dark") else NAVY, bold=True)
    rule(slide, cx - 0.45, 4.30, 0.9, GOLD, weight=3.0)
    if rec.get("subtitle"):
        tb, tf = textbox(slide, 0.8, 4.58, PAGE_W - 1.6, 0.6)
        p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.2)
        run(p, rec["subtitle"], font=SERIF, size=18,
            color=CREAM_DIM if rec.get("dark") else INK, italic=True)
    meta = rec.get("meta")
    if meta:
        tb, tf = textbox(slide, 0, 5.34, PAGE_W, 0.4)
        p = para(tf, first=True, align=PP_ALIGN.CENTER)
        for m in meta:
            run(p, m["label"] + "  ", font=SERIF, size=14, color=GOLD_DK, bold=True, spc=1.2)
            run(p, m["value"] + "        ", font=SERIF, size=14,
                color=CREAM_TXT if rec.get("dark") else NAVY)
    corners(slide, dark=rec.get("dark", False))


def r_divider(slide, rec):
    dark = rec.get("dark", True)
    bg(slide, NAVY_DEEP if dark else CREAM)
    cx = PAGE_W / 2
    cy = rec.get("cy", 2.62)
    if rec.get("eyebrow"):
        _eyebrow(slide, rec["eyebrow"], x=0, y=cy,
                 color=GOLD if dark else GOLD_DK, align=PP_ALIGN.CENTER, w=PAGE_W)
        cy += 0.42
    title = rec.get("title", "")
    tsize = rec.get("title_size", 44)
    # shrink to keep a single clean line where possible (floor 34pt)
    while tsize > 34 and est_lines(title, tsize, PAGE_W - 1.6) > 1:
        tsize -= 2
    tlines = est_lines(title, tsize, PAGE_W - 1.6)
    tb, tf = textbox(slide, 0.6, cy, PAGE_W - 1.2, max(1.4, tlines * 1.0))
    p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.04)
    run(p, title, font=SERIF, size=tsize, color=CREAM_TXT if dark else NAVY, bold=True)
    ry = cy + max(1.18, tlines * (tsize / 72.0) * 1.34) + 0.14
    rule(slide, cx - 0.45, ry, 0.9, GOLD, weight=3.0)
    if rec.get("subtitle"):
        tb, tf = textbox(slide, 0.9, ry + 0.27, PAGE_W - 1.8, 0.9)
        p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.25)
        run(p, rec["subtitle"], font=SERIF, size=rec.get("sub_size", 17),
            color=CREAM_DIM if dark else INK, italic=True)
    corners(slide, dark=dark)


def r_quote(slide, rec):
    dark = rec.get("dark", False)
    bg(slide, NAVY_DEEP if dark else CREAM)
    tb, tf = textbox(slide, MARGIN_L - 0.15, 1.30, 3.0, 2.0)
    p = para(tf, first=True)
    run(p, "“", font=SERIF, size=150, color=GOLD, bold=True)
    qsize = rec.get("title_size", 27)
    tb, tf = textbox(slide, MARGIN_L, 2.95, CONTENT_W, 2.6)
    p = para(tf, first=True, line=1.22)
    run(p, rec.get("quote", rec.get("title", "")), font=SERIF, size=qsize,
        color=CREAM_TXT if dark else NAVY)
    yb = rec.get("attr_y", 5.55)
    rule(slide, MARGIN_L, yb, 0.5, GOLD, weight=2.8)
    if rec.get("attribution"):
        _eyebrow(slide, rec["attribution"], x=MARGIN_L, y=yb + 0.17,
                 color=CREAM_DIM if dark else GOLDM)
    corners(slide, dark=dark)


def r_two_col(slide, rec):
    bg(slide, CREAM)
    if rec.get("ghost"):
        # subtle watermark numeral behind the right column, below the title line
        # (clear of long titles up top and the br corner down low)
        tb, tf = textbox(slide, CONTENT_R - 3.2, 2.95, 3.2, 2.0)
        p = para(tf, first=True, align=PP_ALIGN.RIGHT)
        run(p, rec["ghost"], font=SERIF, size=rec.get("ghost_size", 150),
            color=GHOST, bold=True)
    y = _title_block(slide, rec, title_size=rec.get("title_size", 32))
    cy = y + 0.15
    L = rec["left"]; R = rec["right"]
    c1x, c1w = MARGIN_L, 4.3
    gap_x = 0.65
    c2x = MARGIN_L + c1w + gap_x
    c2w = CONTENT_R - c2x
    rect(slide, c2x - 0.34, cy + 0.05, 0.012, 2.7, LINE)
    # left
    tb, tf = textbox(slide, c1x, cy, c1w, 0.3)
    p = para(tf, first=True)
    run(p, L.get("label", "").upper(), font=SERIF, size=12.5, color=GOLD_DK, bold=True, spc=2.2)
    if L.get("bullets"):
        _bullets(slide, L["bullets"], c1x, cy + 0.44, c1w, 3.0, size=rec.get("body_size", 15.5))
    else:
        tb, tf = textbox(slide, c1x, cy + 0.44, c1w, 3.0)
        p = para(tf, first=True, line=1.3)
        run(p, L.get("body", ""), font=SERIF, size=rec.get("body_size", 16), color=INK)
    # right
    tb, tf = textbox(slide, c2x, cy, c2w, 0.3)
    p = para(tf, first=True)
    run(p, R.get("label", "").upper(), font=SERIF, size=12.5, color=GOLD_DK, bold=True, spc=2.2)
    if R.get("bullets"):
        _bullets(slide, R["bullets"], c2x, cy + 0.44, c2w, 3.2, size=rec.get("body_size", 15.5))
    else:
        tb, tf = textbox(slide, c2x, cy + 0.44, c2w, 3.2)
        p = para(tf, first=True, line=1.3)
        run(p, R.get("body", ""), font=SERIF, size=rec.get("body_size", 16), color=INK)
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_br"))


def r_numbered(slide, rec):
    bg(slide, CREAM)
    y = _title_block(slide, rec, title_size=rec.get("title_size", 32))
    rows = rec["rows"]
    n = len(rows)
    avail = (PAGE_H - MARGIN_B - 0.35) - (y + 0.05)
    rowh = min(rec.get("rowh", 1.12), avail / n)
    yy = y + 0.08
    body_size = rec.get("body_size", 13.5)
    # adaptive num column: widen + shrink for long labels (e.g. times "10:35")
    maxlen = max((len(str(r.get("num", ""))) for r in rows), default=2)
    if maxlen <= 2:
        num_sz, num_w = 26, 0.85
    elif maxlen <= 4:
        num_sz, num_w = 21, 1.25
    else:
        num_sz, num_w = 18, 1.5
    lead_x = MARGIN_L + num_w + 0.22
    for i, row in enumerate(rows):
        if i > 0:
            rect(slide, MARGIN_L, yy - 0.09, CONTENT_W, 0.012, LINE)
        tb, tf = textbox(slide, MARGIN_L, yy - 0.02, num_w, 0.8)
        p = para(tf, first=True)
        run(p, row.get("num", f"{i+1:02d}"), font=SERIF, size=num_sz, color=GOLD_DK, bold=True)
        tb, tf = textbox(slide, lead_x, yy, CONTENT_R - lead_x, rowh)
        p = para(tf, first=True, after=3)
        run(p, row.get("lead", ""), font=SERIF, size=rec.get("lead_size", 17), color=NAVY, bold=True)
        if row.get("body"):
            p = para(tf, line=1.2)
            run(p, row["body"], font=SERIF, size=body_size, color=INK)
        yy += rowh
    if rec.get("kicker"):
        tb, tf = textbox(slide, MARGIN_L, yy + 0.05, CONTENT_W, 0.4)
        p = para(tf, first=True)
        run(p, rec["kicker"], font=SERIF, size=14.5, color=TERRA, italic=True, bold=True)
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_br"))


def r_bullets(slide, rec):
    bg(slide, CREAM)
    y = _title_block(slide, rec, title_size=rec.get("title_size", 32))
    cy = y + 0.12
    if rec.get("intro"):
        tb, tf = textbox(slide, MARGIN_L, cy, CONTENT_W, 0.8)
        p = para(tf, first=True, line=1.25)
        run(p, rec["intro"], font=SERIF, size=rec.get("body_size", 16.5), color=INK)
        cy += 0.30 + 0.26 * (1 + len(rec["intro"]) // 95)
    items = rec["bullets"]
    size = rec.get("body_size", 16.5)
    two = rec.get("two_col_bullets") or (len(items) >= 6 and rec.get("two_col_bullets") is not False)
    if two and len(items) >= 4:
        half = (len(items) + 1) // 2
        colw = (CONTENT_W - 0.6) / 2
        _bullets(slide, items[:half], MARGIN_L, cy, colw, 3.2, size=size, gap=rec.get("gap", 10))
        _bullets(slide, items[half:], MARGIN_L + colw + 0.6, cy, colw, 3.2, size=size, gap=rec.get("gap", 10))
    else:
        _bullets(slide, items, MARGIN_L, cy, CONTENT_W, 3.4, size=size, gap=rec.get("gap", 11))
    if rec.get("kicker"):
        # bottom-anchored, ending above the content margin so it clears the footer
        tb, tf = textbox(slide, MARGIN_L, PAGE_H - MARGIN_B - 1.05, CONTENT_W, 0.92,
                         anchor=MSO_ANCHOR.BOTTOM)
        p = para(tf, first=True, line=1.16)
        run(p, rec["kicker"], font=SERIF, size=14.5, color=TERRA, italic=True, bold=True)
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_br"))


def r_bars(slide, rec):
    bg(slide, CREAM)
    y = _title_block(slide, rec, title_size=rec.get("title_size", 31))
    bars = rec["bars"]
    n = len(bars)
    yy = y + 0.10
    avail = (PAGE_H - MARGIN_B - 0.2) - yy
    gap = 0.1
    barh = min(rec.get("barh", 0.74), (avail - gap * (n - 1)) / n)
    for b in bars:
        fill = _col(b.get("color", "navy"))
        txt = _col(b.get("txt", "cream"), CREAM_TXT)
        rect(slide, MARGIN_L, yy, CONTENT_W, barh, fill)
        tb, tf = textbox(slide, MARGIN_L + 0.28, yy, 0.85, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True); run(p, b.get("num", ""), font=SERIF, size=20, color=txt, bold=True)
        tb, tf = textbox(slide, MARGIN_L + 1.15, yy, 2.55, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True); run(p, b.get("name", ""), font=SERIF, size=18, color=txt, bold=True, spc=1.0)
        tb, tf = textbox(slide, MARGIN_L + 3.8, yy, CONTENT_W - 4.0, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True, line=1.05); run(p, b.get("desc", ""), font=SERIF, size=13, color=txt)
        yy += barh + gap
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_tr"))


def r_exercise(slide, rec):
    bg(slide, CREAM)
    y = _title_block(slide, rec, title_size=rec.get("title_size", 30))
    cx = MARGIN_L
    cy = y + 0.10
    card_lines = rec.get("card", [])
    total_chars = sum(len(x) for x in card_lines)
    # a long card claims full width and drops the side visual
    want_visual = rec.get("visual") == "laptop" or rec.get("image")
    has_visual = bool(want_visual) and total_chars <= 190
    cw = 4.95 if has_visual else CONTENT_W
    csize = rec.get("card_size", 17)
    if total_chars > 300:
        csize = 15
    if total_chars > 460:
        csize = 13.5
    inner_w = cw - 0.8
    nlines = sum(est_lines(x, csize, inner_w) for x in card_lines)
    line_h = (csize / 72.0) * 1.34
    ch = 0.30 * 2 + nlines * line_h + max(0, len(card_lines) - 1) * 0.05
    q = rec.get("question")
    q_reserve = 0.85 if q else 0.0
    avail = (PAGE_H - MARGIN_B - 0.22) - cy - q_reserve
    if not has_visual:
        ch = min(ch, avail)
    else:
        ch = min(ch, PAGE_H - MARGIN_B - 0.4 - cy)
    ch = max(ch, 1.0)
    rect(slide, cx, cy, cw, ch, PAPER, line=LINE, line_w=1.25)
    rect(slide, cx, cy, 0.09, ch, GOLD)
    tb, tf = textbox(slide, cx + 0.42, cy + 0.26, cw - 0.8, ch - 0.4)
    for i, ln in enumerate(card_lines):
        p = para(tf, first=(i == 0), line=1.24, after=5)
        run(p, ln, font=SERIF, size=csize, color=NAVY, bold=True)
    if q:
        qy = cy + ch + 0.20
        tb, tf = textbox(slide, cx, qy, (cw + 1.4 if has_visual else CONTENT_W), 0.9)
        p = para(tf, first=True, line=1.24)
        run(p, "→  ", font=SERIF, size=16, color=TERRA, bold=True)
        run(p, q, font=SERIF, size=15.5, color=INK, italic=True)
    if has_visual and rec.get("image"):
        place_image(slide, rec["image"], 7.7, cy, 3.65, 2.45, fit="contain", frame=True)
    elif has_visual and rec.get("visual") == "laptop":
        _laptop(slide, 7.65, cy + 0.2)
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_br"))


def _laptop(slide, lx, ly, lw=3.7, lh=2.3):
    rect(slide, lx, ly, lw, lh, NAVY)
    rect(slide, lx + 0.13, ly + 0.13, lw - 0.26, lh - 0.26, CREAM)
    rect(slide, lx - 0.42, ly + lh, lw + 0.84, 0.15, NAVY)
    sx, sy = lx + 0.45, ly + 0.5
    for w in (1.8, 2.3, 1.25):
        rect(slide, sx, sy, w, 0.09, LINE); sy += 0.32
    rect(slide, sx, sy, 0.09, 0.22, GOLD)


def r_image_text(slide, rec):
    bg(slide, CREAM)
    side = rec.get("image_side", "right")
    y = _title_block(slide, rec, title_size=rec.get("title_size", 30))
    cy = y + 0.12
    img_w = rec.get("img_w", 4.6)
    text_w = CONTENT_W - img_w - 0.6
    if side == "right":
        tx = MARGIN_L; ix = CONTENT_R - img_w
    else:
        ix = MARGIN_L; tx = MARGIN_L + img_w + 0.6
    img_h = rec.get("img_h", 3.4)
    if rec.get("image"):
        place_image(slide, rec["image"], ix, cy, img_w, img_h,
                    fit=rec.get("fit", "contain"), frame=rec.get("frame", True))
    if rec.get("bullets"):
        _bullets(slide, rec["bullets"], tx, cy, text_w, img_h, size=rec.get("body_size", 16))
    elif rec.get("body"):
        tb, tf = textbox(slide, tx, cy, text_w, img_h)
        p = para(tf, first=True, line=1.3)
        run(p, rec["body"], font=SERIF, size=rec.get("body_size", 16.5), color=INK)
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_br"))


def r_image_full(slide, rec):
    dark = rec.get("dark", True)
    bg(slide, NAVY_DEEP if dark else CREAM)
    # framed image, near full bleed within margins
    bx, by = 0.55, 0.55
    bw, bh = PAGE_W - 1.1, PAGE_H - 1.1
    if rec.get("eyebrow") or rec.get("caption"):
        bh -= 0.55
        by = 0.5
    if rec.get("image"):
        place_image(slide, rec["image"], bx, by, bw, bh,
                    fit=rec.get("fit", "contain"), frame=rec.get("frame", True),
                    frame_color=NAVY if not dark else None)
    if rec.get("eyebrow"):
        _eyebrow(slide, rec["eyebrow"], x=bx + 0.05, y=0.18,
                 color=GOLD if dark else GOLD_DK)
    if rec.get("caption"):
        tb, tf = textbox(slide, bx, PAGE_H - 0.66, bw, 0.46)
        p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.15)
        run(p, rec["caption"], font=SERIF, size=14.5,
            color=CREAM_TXT if dark else INK, italic=True)
    # full-bleed images stay clean (no corner motif over the screenshot)
    corners(slide, dark=dark, **_corner_kw(rec, default="none", drop_dark=True))


def r_flow(slide, rec):
    """Horizontal node/workflow diagram: connected labelled stage boxes
    (e.g. Trigger → Inputs → Instructions → Model → Output)."""
    bg(slide, CREAM)
    y = _title_block(slide, rec, title_size=rec.get("title_size", 30))
    steps = rec["steps"]
    n = len(steps)
    palette = [NAVY, TERRA, SPICE, GOLD, GREEN, NAVY]
    txt_for = [CREAM_TXT, CREAM_TXT, CREAM_TXT, NAVY, CREAM_TXT, CREAM_TXT]
    arrow = 0.42
    avail_h = (PAGE_H - MARGIN_B - 0.35) - (y + 0.1)
    bw = (CONTENT_W - arrow * (n - 1)) / n
    bh = min(rec.get("box_h", 1.7), avail_h - (0.8 if rec.get("caption") else 0))
    cy = y + 0.1 + (avail_h - bh) / 2
    lab_sz = 16 if n <= 4 else (14 if n <= 5 else 12.5)
    sub_sz = 11.5 if n <= 5 else 10.5
    x = MARGIN_L
    for i, step in enumerate(steps):
        fill = palette[i % len(palette)]
        tc = txt_for[i % len(txt_for)]
        rect(slide, x, cy, bw, bh, fill)
        label = step.get("label", "") if isinstance(step, dict) else str(step)
        sub = step.get("sub", "") if isinstance(step, dict) else ""
        tb, tf = textbox(slide, x + 0.12, cy + 0.12, bw - 0.24, bh - 0.24,
                         anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.04)
        run(p, label, font=SERIF, size=lab_sz, color=tc, bold=True)
        if sub:
            p = para(tf, align=PP_ALIGN.CENTER, before=3, line=1.06)
            run(p, sub, font=SERIF, size=sub_sz, color=tc)
        if i < n - 1:
            ax = x + bw
            tb, tf = textbox(slide, ax, cy, arrow, bh, anchor=MSO_ANCHOR.MIDDLE)
            p = para(tf, first=True, align=PP_ALIGN.CENTER)
            run(p, "→", font=SERIF, size=20, color=GOLD_DK, bold=True)
        x += bw + arrow
    if rec.get("caption"):
        tb, tf = textbox(slide, MARGIN_L, cy + bh + 0.22, CONTENT_W, 0.7)
        p = para(tf, first=True, line=1.2)
        run(p, rec["caption"], font=SERIF, size=14, color=INK, italic=True)
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_tr"))


def r_team(slide, rec):
    bg(slide, CREAM)
    y = _title_block(slide, rec, title_size=rec.get("title_size", 31))
    people = rec["people"]
    n = len(people)
    cy = y + 0.18
    pw = rec.get("photo_w", 2.35)
    gap = (CONTENT_W - n * pw) / max(1, n - 1) if n > 1 else 0
    gap = min(gap, 0.9)
    total = n * pw + (n - 1) * gap
    startx = MARGIN_L + (CONTENT_W - total) / 2
    ph = rec.get("photo_h", 2.55)
    for i, person in enumerate(people):
        x = startx + i * (pw + gap)
        if person.get("photo"):
            place_image(slide, person["photo"], x, cy, pw, ph, fit="cover",
                        frame=True, frame_color=GOLD_DK)
        ty = cy + ph + 0.16
        tb, tf = textbox(slide, x - 0.25, ty, pw + 0.5, 1.4)
        p = para(tf, first=True, align=PP_ALIGN.CENTER)
        run(p, person.get("name", ""), font=SERIF, size=17, color=NAVY, bold=True)
        if person.get("role"):
            p = para(tf, align=PP_ALIGN.CENTER, before=2)
            run(p, person["role"], font=SERIF, size=12.5, color=TERRA, italic=True)
        if person.get("sub"):
            p = para(tf, align=PP_ALIGN.CENTER, before=1)
            run(p, person["sub"], font=SERIF, size=11.5, color=GOLDM)
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_tr"))


def r_image_grid(slide, rec):
    bg(slide, CREAM)
    y = _title_block(slide, rec, title_size=rec.get("title_size", 31))
    imgs = rec["images"]
    cy = y + 0.12
    n = len(imgs)
    cols = rec.get("cols", min(n, 4 if n > 3 else n))
    rows = (n + cols - 1) // cols
    gap = 0.28
    gw = (CONTENT_W - gap * (cols - 1)) / cols
    avail_h = (PAGE_H - MARGIN_B - 0.2) - cy
    gh = min(rec.get("cell_h", (avail_h - gap * (rows - 1)) / rows),
             (avail_h - gap * (rows - 1)) / rows)
    for i, im in enumerate(imgs):
        r, c = divmod(i, cols)
        x = MARGIN_L + c * (gw + gap)
        yy = cy + r * (gh + gap)
        fn = im if isinstance(im, str) else im.get("file")
        place_image(slide, fn, x, yy, gw, gh, fit=rec.get("fit", "cover"), frame=True)
    _maybe_footer(slide, rec)
    corners(slide, **_corner_kw(rec, default="tl_tr"))


# ---------------------------------------------------------------- helpers
def _maybe_footer(slide, rec):
    if rec.get("page") is not None:
        footer(slide, page=rec.get("page"), dark=rec.get("dark", False))


def _corner_kw(rec, default="tl_br", drop_dark=False):
    """Translate rec['corners'] into kwargs for corners()."""
    which = rec.get("corners", default)
    mp = {
        "tl_br": ("tl", "br"), "tl_tr": ("tl", "tr"),
        "bl_br": ("bl", "br"), "all": ("tl", "tr", "bl", "br"),
        "tl": ("tl",), "none": (),
    }
    sel = mp.get(which, ("tl", "br"))
    kw = {"which": sel}
    if not drop_dark:
        kw["dark"] = rec.get("dark", False)
    if not sel:
        # render nothing: pass empty which
        kw["which"] = ()
    return kw


RENDERERS = {
    "cover": r_cover,
    "divider": r_divider,
    "statement": r_divider,        # statement == divider on cream (dark=False)
    "quote": r_quote,
    "two_col": r_two_col,
    "comparison": r_two_col,
    "numbered": r_numbered,
    "bullets": r_bullets,
    "bars": r_bars,
    "exercise": r_exercise,
    "image_text": r_image_text,
    "image_full": r_image_full,
    "image_grid": r_image_grid,
    "team": r_team,
    "flow": r_flow,
}


def render(slide, rec):
    arch = rec.get("arch", "bullets")
    fn = RENDERERS.get(arch)
    if fn is None:
        raise ValueError(f"unknown archetype: {arch}")
    if arch == "image_full":
        fn(slide, rec)
    else:
        fn(slide, rec)
