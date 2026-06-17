"""Build the full redesigned deck from archetype records (deck_records.json)."""
import json
import sys
from pptx import Presentation
from theme_indic import EMU_W, EMU_H, ensure_assets
from archetypes import render

FOOTER_ARCH = {"two_col", "comparison", "numbered", "bullets", "bars",
               "exercise", "image_text", "image_grid"}
NO_FOOTER = {"cover", "divider", "statement", "quote", "image_full"}


def build(records_path="deck_records.json", out="IAF-Session1-Redesigned.pptx"):
    ensure_assets()
    records = json.load(open(records_path))
    if isinstance(records, dict):
        records = records.get("records", [])

    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    BLANK = prs.slide_layouts[6]

    page = 0
    for idx, rec in enumerate(records, 1):
        arch = rec.get("arch", "bullets")
        # sequential footer page number on content archetypes
        if arch in FOOTER_ARCH:
            page += 1
            rec["page"] = page
        else:
            rec.pop("page", None)
        slide = prs.slides.add_slide(BLANK)
        try:
            render(slide, rec)
        except Exception as e:
            print(f"  !! render error on record {idx} src={rec.get('src')} arch={arch}: {e}")
            raise
        # speaker note for flagged images
        note = rec.get("note")
        if note:
            slide.notes_slide.notes_text_frame.text = note
        # preserve hidden status of source slide 4
        if str(rec.get("src", "")).rstrip("ab") == "4" and arch != "cover":
            if str(rec.get("src")) == "4":
                slide._element.set("show", "0")

    prs.save(out)
    print(f"saved {out}  ·  {len(prs.slides._sldIdLst)} slides  ·  {page} numbered")
    return out


if __name__ == "__main__":
    rp = sys.argv[1] if len(sys.argv) > 1 else "deck_records.json"
    op = sys.argv[2] if len(sys.argv) > 2 else "IAF-Session1-Redesigned.pptx"
    build(rp, op)
