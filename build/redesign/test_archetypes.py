"""Smoke-test every archetype renderer with representative records."""
import sys
from pptx import Presentation
from theme_indic import EMU_W, EMU_H, ensure_assets
from archetypes import render

ensure_assets()
prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

RECS = [
    {"arch": "cover", "eyebrow": [{"t": "DISCOVERY DAY", "c": "terra"},
        {"t": "    ·    BANGALORE    ·    JUNE 2026", "c": "gold"}],
     "title": "ImpactAI Foundry", "subtitle": "Welcome — find a seat, grab some snacks.",
     "meta": [{"label": "WiFi", "value": "ImpactAI-Guest"}, {"label": "Password", "value": "foundry2026"}]},
    {"arch": "divider", "dark": True, "eyebrow": "SESSION ONE   ·   10:00 AM",
     "title": "Prompting and design", "subtitle": "The shortest path from a blank prompt to useful output."},
    {"arch": "statement", "dark": False, "eyebrow": "BEFORE WE BUILD",
     "title": "You're in control", "subtitle": "The tools serve your mission, not the other way around."},
    {"arch": "quote", "quote": "You know your work better than anyone in this room. Our job isn't to tell you what to build — it's to help you put AI to the work you already do.",
     "attribution": "THE PREMISE OF THE FOUNDRY", "page": 16},
    {"arch": "two_col", "eyebrow": "PHASE 01 · DISCOVER", "title": "Ideation & foundational learning",
     "ghost": "01", "left": {"label": "Focus", "body": "Explore where AI can help and build hands-on fluency before committing to a problem."},
     "right": {"label": "Goals", "bullets": ["Explore high-impact AI use cases", "Build fluency with prompting and core tools", "Understand AI's strengths and limits", "Spot workflows that feel unclear or manual"]}, "page": 20},
    {"arch": "numbered", "eyebrow": "THREE GUARDRAILS", "title": "Build with care, but build",
     "rows": [{"num": "01", "lead": "Protect personal information", "body": "Free AI tools aren't private by default. Beneficiary and donor data stays out unless the platform is secure."},
              {"num": "02", "lead": "Keep a human in the loop", "body": "Until you fully understand how a tool behaves on your work, review its outputs before relying on them."},
              {"num": "03", "lead": "Customise for your community", "body": "Out-of-the-box AI doesn't reflect every context. Adapt it to truly serve the people you serve."}],
     "kicker": "…but don't let these get in the way of getting started.", "page": 30},
    {"arch": "bullets", "eyebrow": "HABITS THAT SAVE TOKENS", "title": "Spend less, get more",
     "bullets": ["Start a fresh chat for each new task", "Be specific — vague prompts waste turns",
                 "Give examples of what good looks like", "Edit and re-run rather than starting over",
                 "Use the cheapest model that does the job", "Keep reference docs in one place"], "page": 34},
    {"arch": "bars", "eyebrow": "THE LEARNING PATH", "title": "Five phases across four sessions",
     "bars": [{"num": "01", "name": "DISCOVER", "desc": "Build prompting fluency; spot AI-shaped problems.", "color": "navy"},
              {"num": "02", "name": "DEFINE", "desc": "Narrow to one workflow; map inputs and success.", "color": "terra"},
              {"num": "03", "name": "DESIGN", "desc": "Build a reusable prototype with your mentor.", "color": "spice"},
              {"num": "04", "name": "REFINE", "desc": "Stress-test edge cases; fit it to your team.", "color": "gold", "txt": "navy"},
              {"num": "05", "name": "DELIVER", "desc": "Demo Day, August 1 — share what you built.", "color": "green"}], "page": 19},
    {"arch": "exercise", "eyebrow": "EXERCISE · 15 MIN", "title": "Ask AI to solve a real work problem",
     "card": ["You're a 10th-grade English teacher. You just hired a teaching assistant to grade essays."],
     "question": "With your team: what do you do to make sure they do the right thing?", "visual": "laptop", "page": 36},
    {"arch": "image_text", "eyebrow": "SOLUTION TYPE · 01", "title": "The Assistant",
     "image": "image41.png", "image_side": "right",
     "bullets": ["Drafts and edits on demand", "Answers questions in plain language", "Best for one-off, ad-hoc help"], "page": 49},
    {"arch": "image_full", "dark": True, "image": "image16.png", "caption": "Claude — a working conversation", "eyebrow": "A LIVE LOOK"},
    {"arch": "image_grid", "eyebrow": "MEET YOUR COHORT", "title": "Twelve organisations, one room",
     "images": ["image12.jpg", "image22.jpg", "image24.jpg", "image40.png", "image44.png", "image46.png", "image12.jpg", "image22.jpg"], "cols": 4, "page": 5},
]

for rec in RECS:
    s = prs.slides.add_slide(BLANK)
    render(s, rec)

out = sys.argv[1] if len(sys.argv) > 1 else "test_archetypes.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
