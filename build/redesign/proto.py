"""Prototype: 7 archetype slides in the new warm-minimal indic system."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from theme import *  # noqa

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def new(bgcolor=WHITE):
    s = prs.slides.add_slide(BLANK)
    bg(s, bgcolor)
    return s


def ghost_numeral(slide, text, color, x=None, y=0.45, size=210):
    """Big faint phase numeral, sits behind content as identity mark."""
    if x is None:
        x = CONTENT_R - 3.0
    tb, tf = textbox(slide, x, y, 3.0, 3.0, anchor=MSO_ANCHOR.TOP)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT)
    run(p, text, font=DISPLAY, size=size, color=color, bold=True)


# ===========================================================================
# 1 — COVER
# ===========================================================================
def s_cover():
    s = new(CREAM)
    corner_motif(s, SAFFRON, x=CONTENT_R - 0.5, y=MARGIN_T)
    # eyebrow
    tb, tf = textbox(s, MARGIN_L, 2.30, CONTENT_W, 0.4)
    p = para(tf, first=True)
    run(p, "DISCOVERY DAY", font=EYEBROW_FONT, size=13, color=TERRACOTTA, bold=True, spc=2.6)
    run(p, "   ·   BANGALORE   ·   JUNE 2026", font=EYEBROW_FONT, size=13, color=MUTE, bold=True, spc=2.6)
    # title
    tb, tf = textbox(s, MARGIN_L, 2.74, CONTENT_W, 1.5)
    p = para(tf, first=True, line=1.0)
    run(p, "ImpactAI Foundry", font=DISPLAY, size=68, color=INDIGO, bold=True)
    rule(s, MARGIN_L, 4.18, 0.9, SAFFRON, weight=4.5)
    # subtitle
    tb, tf = textbox(s, MARGIN_L, 4.46, 9.0, 0.6)
    p = para(tf, first=True)
    run(p, "Welcome — find a seat, grab some snacks.", font=BODY, size=18, color=INK)
    # wifi pills
    px, py = MARGIN_L, 5.30
    for label, val in [("WiFi", "ImpactAI-Guest"), ("Password", "foundry2026")]:
        chip = rect(s, px, py, 0.0, 0.0, None)  # placeholder; sized via text box widths
        # label box
        w_total = 3.05
        r = rect(s, px, py, w_total, 0.52, WHITE, line=LINE, line_w=1.0)
        tbx, tfx = textbox(s, px + 0.22, py, w_total - 0.4, 0.52, anchor=MSO_ANCHOR.MIDDLE)
        pp = para(tfx, first=True)
        run(pp, label.upper() + "   ", font=EYEBROW_FONT, size=10, color=MUTE, bold=True, spc=1.4)
        run(pp, val, font=DISPLAY, size=14, color=INDIGO, bold=True)
        px += w_total + 0.35
    footer(s, page=None)
    return s


# ===========================================================================
# 2 — SECTION DIVIDER (indigo)
# ===========================================================================
def s_divider():
    s = new(INDIGO_DEEP)
    # subtle motif
    corner_motif(s, SAFFRON, x=CONTENT_R - 0.5, y=MARGIN_T)
    rule(s, MARGIN_L, 2.55, 0.9, SAFFRON, weight=4.5)
    tb, tf = textbox(s, MARGIN_L, 2.82, CONTENT_W, 0.4)
    p = para(tf, first=True)
    run(p, "SESSION ONE   ·   10:00 AM", font=EYEBROW_FONT, size=13, color=SAFFRON, bold=True, spc=2.6)
    tb, tf = textbox(s, MARGIN_L, 3.20, CONTENT_W, 1.4)
    p = para(tf, first=True, line=1.0)
    run(p, "Prompting and design", font=DISPLAY, size=54, color=WHITE, bold=True)
    tb, tf = textbox(s, MARGIN_L, 4.55, 9.5, 0.6)
    p = para(tf, first=True)
    run(p, "The shortest path from a blank prompt to useful output.",
        font=BODY, size=18, color=CREAM_DIM)
    footer(s, page=24, dark=True)
    return s


# ===========================================================================
# 3 — TWO-COLUMN PHASE
# ===========================================================================
def s_phase():
    s = new(WHITE)
    ghost_numeral(s, "01", RGBColor(0xEE, 0xF0, 0xF6), x=CONTENT_R - 3.0, y=0.30, size=210)
    eyebrow(s, "PHASE 01 · DISCOVER")
    tb, tf = textbox(s, MARGIN_L, 1.18, 9.0, 1.0)
    p = para(tf, first=True, line=1.0)
    run(p, "Ideation & foundational learning", font=DISPLAY, size=34, color=INDIGO, bold=True)
    rule(s, MARGIN_L, 1.92, 0.62, SAFFRON, weight=3.5)
    # columns
    cy = 2.55
    col1_x, col1_w = MARGIN_L, 4.7
    col2_x, col2_w = 6.4, CONTENT_R - 6.4
    # col1
    tb, tf = textbox(s, col1_x, cy, col1_w, 0.3)
    p = para(tf, first=True); run(p, "FOCUS", font=EYEBROW_FONT, size=12, color=SAFFRON_DK, bold=True, spc=2.0)
    tb, tf = textbox(s, col1_x, cy + 0.42, col1_w, 2.4)
    p = para(tf, first=True, line=1.28)
    run(p, "Explore where AI can help and build hands-on fluency before committing to a problem.",
        font=BODY, size=17, color=INK)
    # col2
    tb, tf = textbox(s, col2_x, cy, col2_w, 0.3)
    p = para(tf, first=True); run(p, "GOALS", font=EYEBROW_FONT, size=12, color=SAFFRON_DK, bold=True, spc=2.0)
    goals = [
        "Explore high-impact AI use cases",
        "Build fluency with prompting and core tools",
        "Understand AI's strengths and limits",
        "Spot workflows that feel unclear or manual",
    ]
    tb, tf = textbox(s, col2_x, cy + 0.42, col2_w, 3.4)
    for i, g in enumerate(goals):
        p = para(tf, first=(i == 0), line=1.18, after=8)
        run(p, "—  ", font=BODY, size=17, color=SAFFRON, bold=True)
        run(p, g, font=BODY, size=17, color=INK)
    footer(s, page=20)
    return s


# ===========================================================================
# 4 — NUMBERED CONTENT
# ===========================================================================
def s_numbered():
    s = new(WHITE)
    eyebrow(s, "THREE GUARDRAILS")
    tb, tf = textbox(s, MARGIN_L, 1.18, 10.0, 1.0)
    p = para(tf, first=True, line=1.0)
    run(p, "Build with care, but build", font=DISPLAY, size=34, color=INDIGO, bold=True)
    rule(s, MARGIN_L, 1.92, 0.62, SAFFRON, weight=3.5)
    rows = [
        ("01", "Protect personal information", "Free AI tools aren't private by default. Beneficiary and donor data stays out unless the platform is secure and configured for it."),
        ("02", "Keep a human in the loop", "Until you fully understand how a tool behaves on your work, review its outputs before relying on them."),
        ("03", "Customise for your community", "Out-of-the-box AI doesn't reflect every context. Adapt it to truly serve the people you serve."),
    ]
    y = 2.50
    rowh = 1.18
    for i, (num, lead, body) in enumerate(rows):
        if i > 0:
            rect(s, MARGIN_L, y - 0.10, CONTENT_W, 0.012, LINE)
        # number
        tb, tf = textbox(s, MARGIN_L, y, 0.95, 0.8)
        p = para(tf, first=True)
        run(p, num, font=DISPLAY, size=30, color=SAFFRON, bold=True)
        # lead + body
        tb, tf = textbox(s, MARGIN_L + 1.05, y + 0.02, CONTENT_W - 1.05, rowh)
        p = para(tf, first=True, after=3)
        run(p, lead, font=DISPLAY, size=17.5, color=INDIGO, bold=True)
        p = para(tf, line=1.18)
        run(p, body, font=BODY, size=14.5, color=INK)
        y += rowh
    # closer
    tb, tf = textbox(s, MARGIN_L, y + 0.04, CONTENT_W, 0.5)
    p = para(tf, first=True)
    run(p, "…but don't let these get in the way of getting started.",
        font=BODY, size=15, color=SAND, italic=True, bold=True)
    footer(s, page=30)
    return s


# ===========================================================================
# 5 — EXERCISE (accent band + prompt card + device)
# ===========================================================================
def s_exercise():
    s = new(WHITE)
    # accent band
    band_h = 1.12
    rect(s, 0, 0, PAGE_W, band_h, SAFFRON)
    tb, tf = textbox(s, MARGIN_L, 0, 9.5, band_h, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, after=2)
    run(p, "EXERCISE · 15 MIN", font=EYEBROW_FONT, size=12.5, color=INDIGO, bold=True, spc=2.4)
    p = para(tf, line=1.0)
    run(p, "Ask AI to solve a real work problem", font=DISPLAY, size=28, color=INDIGO, bold=True)
    # prompt card (left)
    cx, cy, cw, ch = MARGIN_L, 1.72, 5.7, 2.45
    rect(s, cx, cy, cw, ch, CREAM, line=LINE, line_w=1.25)
    rect(s, cx, cy, 0.10, ch, SAFFRON)  # left spine
    tb, tf = textbox(s, cx + 0.42, cy + 0.34, cw - 0.7, ch - 0.6)
    p = para(tf, first=True, line=1.22, after=6)
    run(p, "You're a 10th-grade English teacher. You just hired a teaching assistant to grade essays.",
        font=DISPLAY, size=18, color=INDIGO, bold=True)
    # bullets below
    tb, tf = textbox(s, cx, cy + ch + 0.28, cw + 0.4, 1.2)
    p = para(tf, first=True, line=1.2)
    run(p, "→  ", font=BODY, size=16, color=SAFFRON, bold=True)
    run(p, "Work with your team: what do you do to make sure they ", font=BODY, size=16, color=INK)
    run(p, "do the right thing", font=BODY, size=16, color=INK, italic=True)
    run(p, "?", font=BODY, size=16, color=INK)
    # device (right) — minimal flat laptop
    lx, ly, lw, lh = 7.55, 2.05, 4.1, 2.55
    rect(s, lx, ly, lw, lh, INDIGO)               # lid
    rect(s, lx + 0.14, ly + 0.14, lw - 0.28, lh - 0.28, WHITE)  # screen
    rect(s, lx - 0.45, ly + lh, lw + 0.9, 0.16, INDIGO)         # base
    # faint prompt lines on screen
    sx, sy = lx + 0.5, ly + 0.6
    for w in (2.0, 2.6, 1.4):
        rect(s, sx, sy, w, 0.10, LINE); sy += 0.34
    rect(s, sx, sy, 0.10, 0.24, SAFFRON)  # cursor
    footer(s, page=36)
    return s


# ===========================================================================
# 6 — FEATURED QUOTE
# ===========================================================================
def s_quote():
    s = new(CREAM)
    # big quotation mark
    tb, tf = textbox(s, MARGIN_L - 0.12, 0.95, 3.0, 2.0)
    p = para(tf, first=True)
    run(p, "“", font=DISPLAY, size=150, color=SAFFRON, bold=True)
    # quote
    tb, tf = textbox(s, MARGIN_L, 2.55, 10.6, 2.6)
    p = para(tf, first=True, line=1.18)
    run(p, "You know your work better than anyone in this room. Our job isn't to tell you what to build — it's to help you put AI to the work you already do.",
        font=DISPLAY, size=30, color=INDIGO, bold=False)
    # attribution
    rule(s, MARGIN_L, 5.55, 0.5, SAFFRON, weight=3.5)
    tb, tf = textbox(s, MARGIN_L, 5.72, 9.0, 0.4)
    p = para(tf, first=True)
    run(p, "THE PREMISE OF THE FOUNDRY", font=EYEBROW_FONT, size=12, color=MUTE, bold=True, spc=2.2)
    footer(s, page=16)
    return s


# ===========================================================================
# 7 — LEARNING PATH (color bars)
# ===========================================================================
def s_path():
    s = new(WHITE)
    eyebrow(s, "THE LEARNING PATH")
    tb, tf = textbox(s, MARGIN_L, 1.16, 11.0, 0.8)
    p = para(tf, first=True, line=1.0)
    run(p, "Five phases across four sessions", font=DISPLAY, size=32, color=INDIGO, bold=True)
    rule(s, MARGIN_L, 1.86, 0.62, SAFFRON, weight=3.5)
    phases = [
        ("01", "DISCOVER", "Build prompting fluency; spot AI-shaped problems in your work.", INDIGO, WHITE),
        ("02", "DEFINE",   "Narrow to one workflow; map inputs, steps, and success.",       TERRACOTTA, WHITE),
        ("03", "DESIGN",   "Build a reusable prototype with your mentor; test on real work.", SPICE, WHITE),
        ("04", "REFINE",   "Stress-test edge cases; add validation; fit it to your team.",    SAFFRON, INDIGO),
        ("05", "DELIVER",  "Demo Day, August 1 — share what you built and what's next.",      SAND, WHITE),
    ]
    y = 2.36
    barh = 0.82
    gap = 0.09
    for num, name, desc, fill, txt in phases:
        rect(s, MARGIN_L, y, CONTENT_W, barh, fill)
        # number
        tb, tf = textbox(s, MARGIN_L + 0.30, y, 0.9, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True)
        run(p, num, font=DISPLAY, size=22, color=txt, bold=True)
        # name
        tb, tf = textbox(s, MARGIN_L + 1.25, y, 2.6, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True)
        run(p, name, font=DISPLAY, size=20, color=txt, bold=True, spc=0.5)
        # descriptor
        tb, tf = textbox(s, MARGIN_L + 4.0, y, CONTENT_W - 4.2, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True, line=1.05)
        run(p, desc, font=BODY, size=13.5, color=txt)
        y += barh + gap
    footer(s, page=19)
    return s


order = [s_cover, s_divider, s_phase, s_numbered, s_exercise, s_quote, s_path]
for fn in order:
    fn()

out = sys.argv[1] if len(sys.argv) > 1 else "prototype.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
