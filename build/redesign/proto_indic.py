"""Prototype v2 — INDIC direction: cream + kilim corners + Georgia serif."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from theme_indic import *  # noqa

ensure_assets()

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def new(bgcolor=CREAM):
    s = prs.slides.add_slide(BLANK)
    bg(s, bgcolor)
    return s


def ghost(slide, text, color, x, y, size=200):
    tb, tf = textbox(slide, x, y, 3.2, 3.2)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT)
    run(p, text, font=SERIF, size=size, color=color, bold=True)


# ============================================================ 1 COVER
def s_cover():
    s = new(CREAM)
    corners(s)
    cx = PAGE_W / 2
    tb, tf = textbox(s, 0, 2.42, PAGE_W, 0.4)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, "DISCOVERY DAY", font=SERIF, size=13.5, color=TERRA, bold=True, spc=3.0)
    run(p, "    ·    BANGALORE    ·    JUNE 2026", font=SERIF, size=13.5, color=GOLDM, bold=True, spc=3.0)
    tb, tf = textbox(s, 0, 2.86, PAGE_W, 1.4)
    p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.0)
    run(p, "ImpactAI Foundry", font=SERIF, size=62, color=NAVY, bold=True)
    rule(s, cx - 0.45, 4.30, 0.9, GOLD, weight=3.0)
    tb, tf = textbox(s, 0, 4.58, PAGE_W, 0.5)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, "Welcome — find a seat, grab some snacks.", font=SERIF, size=18, color=INK, italic=True)
    tb, tf = textbox(s, 0, 5.28, PAGE_W, 0.4)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, "WiFi  ", font=SERIF, size=14, color=GOLD_DK, bold=True, spc=1.2)
    run(p, "ImpactAI-Guest        ", font=SERIF, size=14, color=NAVY)
    run(p, "Password  ", font=SERIF, size=14, color=GOLD_DK, bold=True, spc=1.2)
    run(p, "foundry2026", font=SERIF, size=14, color=NAVY)
    return s


# ============================================================ 2 DIVIDER
def s_divider():
    s = new(NAVY_DEEP)
    corners(s, dark=True)
    cx = PAGE_W / 2
    eyebrow(s, "SESSION ONE   ·   10:00 AM", x=0, y=2.62, color=GOLD, align=PP_ALIGN.CENTER, w=PAGE_W)
    tb, tf = textbox(s, 0, 3.02, PAGE_W, 1.2)
    p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.0)
    run(p, "Prompting and design", font=SERIF, size=46, color=CREAM_TXT, bold=True)
    rule(s, cx - 0.45, 4.28, 0.9, GOLD, weight=3.0)
    tb, tf = textbox(s, 0, 4.55, PAGE_W, 0.5)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, "The shortest path from a blank prompt to useful output.",
        font=SERIF, size=17, color=CREAM_DIM, italic=True)
    return s


# ============================================================ 3 TWO-COLUMN
def s_phase():
    s = new(CREAM)
    corners(s, which=("tl", "br"))
    ghost(s, "01", PAPER, x=CONTENT_R - 3.0, y=0.7, size=180)
    eyebrow(s, "PHASE 01 · DISCOVER")
    tb, tf = textbox(s, MARGIN_L, 1.56, CONTENT_W, 0.9)
    p = para(tf, first=True, line=1.0)
    run(p, "Ideation & foundational learning", font=SERIF, size=33, color=NAVY, bold=True)
    rule(s, MARGIN_L, 2.32, 0.62, GOLD, weight=2.8)
    cy = 2.95
    c1x, c1w = MARGIN_L, 4.3
    c2x, c2w = MARGIN_L + 4.95, CONTENT_R - (MARGIN_L + 4.95)
    rect(s, MARGIN_L + 4.55, cy + 0.05, 0.012, 2.5, LINE)  # divider
    tb, tf = textbox(s, c1x, cy, c1w, 0.3)
    p = para(tf, first=True); run(p, "FOCUS", font=SERIF, size=12.5, color=GOLD_DK, bold=True, spc=2.2)
    tb, tf = textbox(s, c1x, cy + 0.44, c1w, 2.2)
    p = para(tf, first=True, line=1.3)
    run(p, "Explore where AI can help and build hands-on fluency before committing to a problem.",
        font=SERIF, size=16.5, color=INK)
    tb, tf = textbox(s, c2x, cy, c2w, 0.3)
    p = para(tf, first=True); run(p, "GOALS", font=SERIF, size=12.5, color=GOLD_DK, bold=True, spc=2.2)
    goals = ["Explore high-impact AI use cases",
             "Build fluency with prompting and core tools",
             "Understand AI's strengths and limits",
             "Spot workflows that feel unclear or manual"]
    tb, tf = textbox(s, c2x, cy + 0.44, c2w, 3.0)
    for i, g in enumerate(goals):
        p = para(tf, first=(i == 0), line=1.2, after=9)
        run(p, "•  ", font=SERIF, size=16.5, color=TERRA, bold=True)
        run(p, g, font=SERIF, size=16.5, color=INK)
    footer(s, page=20)
    return s


# ============================================================ 4 NUMBERED
def s_numbered():
    s = new(CREAM)
    corners(s, which=("tl", "br"))
    eyebrow(s, "THREE GUARDRAILS")
    tb, tf = textbox(s, MARGIN_L, 1.56, CONTENT_W, 0.9)
    p = para(tf, first=True, line=1.0)
    run(p, "Build with care, but build", font=SERIF, size=33, color=NAVY, bold=True)
    rule(s, MARGIN_L, 2.32, 0.62, GOLD, weight=2.8)
    rows = [("01", "Protect personal information", "Free AI tools aren't private by default. Beneficiary and donor data stays out unless the platform is secure and configured for it."),
            ("02", "Keep a human in the loop", "Until you fully understand how a tool behaves on your work, review its outputs before relying on them."),
            ("03", "Customise for your community", "Out-of-the-box AI doesn't reflect every context. Adapt it to truly serve the people you serve.")]
    y = 2.86
    rowh = 1.12
    for i, (num, lead, body) in enumerate(rows):
        if i > 0:
            rect(s, MARGIN_L, y - 0.10, CONTENT_W, 0.012, LINE)
        tb, tf = textbox(s, MARGIN_L, y - 0.02, 0.9, 0.8)
        p = para(tf, first=True)
        run(p, num, font=SERIF, size=27, color=GOLD_DK, bold=True)
        tb, tf = textbox(s, MARGIN_L + 1.0, y, CONTENT_W - 1.0, rowh)
        p = para(tf, first=True, after=3)
        run(p, lead, font=SERIF, size=17, color=NAVY, bold=True)
        p = para(tf, line=1.2)
        run(p, body, font=SERIF, size=13.5, color=INK)
        y += rowh
    tb, tf = textbox(s, MARGIN_L, y + 0.06, CONTENT_W, 0.4)
    p = para(tf, first=True)
    run(p, "…but don't let these get in the way of getting started.",
        font=SERIF, size=14.5, color=TERRA, italic=True, bold=True)
    footer(s, page=30)
    return s


# ============================================================ 5 EXERCISE
def s_exercise():
    s = new(CREAM)
    corners(s, which=("tl", "br"))
    eyebrow(s, "EXERCISE · 15 MIN")
    tb, tf = textbox(s, MARGIN_L, 1.56, CONTENT_W, 0.9)
    p = para(tf, first=True, line=1.0)
    run(p, "Ask AI to solve a real work problem", font=SERIF, size=30, color=NAVY, bold=True)
    rule(s, MARGIN_L, 2.30, 0.62, GOLD, weight=2.8)
    # prompt card
    cx, cy, cw, ch = MARGIN_L, 2.85, 4.85, 2.3
    rect(s, cx, cy, cw, ch, PAPER, line=LINE, line_w=1.25)
    rect(s, cx, cy, 0.09, ch, GOLD)
    tb, tf = textbox(s, cx + 0.4, cy + 0.34, cw - 0.7, ch - 0.6)
    p = para(tf, first=True, line=1.25, after=6)
    run(p, "You're a 10th-grade English teacher. You just hired a teaching assistant to grade essays.",
        font=SERIF, size=17, color=NAVY, bold=True)
    tb, tf = textbox(s, cx, cy + ch + 0.26, cw + 1.0, 1.0)
    p = para(tf, first=True, line=1.25)
    run(p, "→  ", font=SERIF, size=16, color=TERRA, bold=True)
    run(p, "With your team: what do you do to make sure they ", font=SERIF, size=15.5, color=INK)
    run(p, "do the right thing", font=SERIF, size=15.5, color=INK, italic=True)
    run(p, "?", font=SERIF, size=15.5, color=INK)
    # laptop (navy outline)
    lx, ly, lw, lh = 7.65, 3.1, 3.7, 2.3
    rect(s, lx, ly, lw, lh, NAVY)
    rect(s, lx + 0.13, ly + 0.13, lw - 0.26, lh - 0.26, CREAM)
    rect(s, lx - 0.42, ly + lh, lw + 0.84, 0.15, NAVY)
    sx, sy = lx + 0.45, ly + 0.5
    for w in (1.8, 2.3, 1.25):
        rect(s, sx, sy, w, 0.09, LINE); sy += 0.32
    rect(s, sx, sy, 0.09, 0.22, GOLD)
    footer(s, page=36)
    return s


# ============================================================ 6 QUOTE
def s_quote():
    s = new(CREAM)
    corners(s)
    tb, tf = textbox(s, MARGIN_L - 0.15, 1.35, 3.0, 2.0)
    p = para(tf, first=True)
    run(p, "“", font=SERIF, size=150, color=GOLD, bold=True)
    tb, tf = textbox(s, MARGIN_L, 2.95, CONTENT_W, 2.4)
    p = para(tf, first=True, line=1.22)
    run(p, "You know your work better than anyone in this room. Our job isn't to tell you what to build — it's to help you put AI to the work you already do.",
        font=SERIF, size=27, color=NAVY)
    rule(s, MARGIN_L, 5.55, 0.5, GOLD, weight=2.8)
    eyebrow(s, "THE PREMISE OF THE FOUNDRY", x=MARGIN_L, y=5.72, color=GOLDM)
    footer(s, page=16)
    return s


# ============================================================ 7 LEARNING PATH
def s_path():
    s = new(CREAM)
    corners(s, which=("tl", "tr"))   # bottom-heavy: bars fill base, frame top
    eyebrow(s, "THE LEARNING PATH")
    tb, tf = textbox(s, MARGIN_L, 1.54, CONTENT_W, 0.8)
    p = para(tf, first=True, line=1.0)
    run(p, "Five phases across four sessions", font=SERIF, size=31, color=NAVY, bold=True)
    rule(s, MARGIN_L, 2.26, 0.62, GOLD, weight=2.8)
    SPICE = RGBColor(0x9E, 0x6B, 0x33)
    phases = [("01", "DISCOVER", "Build prompting fluency; spot AI-shaped problems in your work.", NAVY, CREAM_TXT),
              ("02", "DEFINE", "Narrow to one workflow; map inputs, steps, and success.", TERRA, CREAM_TXT),
              ("03", "DESIGN", "Build a reusable prototype with your mentor; test on real work.", SPICE, CREAM_TXT),
              ("04", "REFINE", "Stress-test edge cases; add validation; fit it to your team.", GOLD, NAVY),
              ("05", "DELIVER", "Demo Day, August 1 — share what you built and what's next.", GREEN, CREAM_TXT)]
    y = 2.72
    barh = 0.74
    gap = 0.1
    for num, name, desc, fill, txt in phases:
        rect(s, MARGIN_L, y, CONTENT_W, barh, fill)
        tb, tf = textbox(s, MARGIN_L + 0.28, y, 0.85, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True); run(p, num, font=SERIF, size=20, color=txt, bold=True)
        tb, tf = textbox(s, MARGIN_L + 1.15, y, 2.5, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True); run(p, name, font=SERIF, size=18, color=txt, bold=True, spc=1.0)
        tb, tf = textbox(s, MARGIN_L + 3.75, y, CONTENT_W - 3.95, barh, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, first=True, line=1.05); run(p, desc, font=SERIF, size=13, color=txt)
        y += barh + gap
    footer(s, page=19)
    return s


for fn in [s_cover, s_divider, s_phase, s_numbered, s_exercise, s_quote, s_path]:
    fn()

out = sys.argv[1] if len(sys.argv) > 1 else "prototype_indic.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
