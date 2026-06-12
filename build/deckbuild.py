"""
deckbuild.py — Build the ImpactAI Foundry Session 1 deck by cloning the
4-variant template (impact-ai-foundry-layouts.pptx) once per target slide and
filling each variant's text slots from a content spec.

Design system (from template + 18-slide reference deck):
  Colors : navy 1F3A5F, terracotta B85042, gold C8A04A, muted-gold 8A7A55,
           green 4A6B3A, cream F4ECD8
  Fonts  : Georgia (serif headings/labels), Calibri (sans body)

Variants and their editable text slots (slot keys used by content specs):
  V1 (full border, centered title) — section dividers / titles
      eyebrow  Georgia 13 bold terracotta, centered
      title    Georgia 40 navy,            centered
      subtitle Calibri 16 muted-gold,      centered
  V2 (full border + side masks = horizontal bands) — workhorse content
      eyebrow  Georgia 10 bold terracotta
      title    Georgia 30 bold navy
      body     rich multi-paragraph block (list of paragraph specs)
  V3 (corner motifs) — featured quote / single idea
      eyebrow  Georgia 10 bold terracotta
      quote    Georgia 22 navy
      attrib   Calibri 13 muted-gold
  V4 (vertical side bands) — two-column pairings
      eyebrow  Georgia 10 bold terracotta
      title    Georgia 26 bold navy
      col1_head Georgia 15 bold terracotta   col1_body rich block
      col2_head Georgia 15 bold green         col2_body rich block

Every slide also has a brand label ("IMPACT AI FOUNDRY") kept constant and a
page-number label set automatically from slide position.
"""
import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

TEMPLATE = "/Users/arjun/Desktop/IAF-Session1/impact-ai-foundry-layouts.pptx"

NAVY  = RGBColor(0x1F, 0x3A, 0x5F)
TERRA = RGBColor(0xB8, 0x50, 0x42)
GOLD  = RGBColor(0xC8, 0xA0, 0x4A)
MGOLD = RGBColor(0x8A, 0x7A, 0x55)
GREEN = RGBColor(0x4A, 0x6B, 0x3A)

COLMAP = {"navy": NAVY, "terra": TERRA, "gold": GOLD, "mgold": MGOLD, "green": GREEN}

# ---------------------------------------------------------------- duplication

def _remap_rels(src_slide, new_slide):
    """After deep-copying XML, re-point every r:embed/r:id/r:link in the new
    slide to the same target parts the source referenced, with fresh rIds."""
    src_part, new_part = src_slide.part, new_slide.part
    # gather rId-bearing attributes in the copied element tree
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    attrs = [qn("r:embed"), qn("r:id"), qn("r:link"), qn("r:pict"),
             qn("r:dm"), qn("r:lo"), qn("r:qs"), qn("r:cs")]
    used = {}
    for el in new_slide._element.iter():
        for a in attrs:
            rid = el.get(a)
            if rid:
                used.setdefault(rid, []).append((el, a))
    for old_rid, hits in used.items():
        if old_rid not in src_part.rels:
            continue
        rel = src_part.rels[old_rid]
        if rel.is_external:
            new_rid = new_part.relate_to(rel.target_ref, rel.reltype,
                                         is_external=True)
        else:
            new_rid = new_part.relate_to(rel.target_part, rel.reltype)
        for el, a in hits:
            el.set(a, new_rid)


def duplicate_slide(prs, src_slide):
    """Create a new slide that is a faithful clone of src_slide, including its
    background blipFill image and all shapes."""
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # strip the placeholders add_slide() inserted from the layout
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    new_cSld = new_slide._element.find(qn("p:cSld"))
    src_cSld = src_slide._element.find(qn("p:cSld"))
    # copy background (<p:bg>) if present — insert as first child of cSld
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        new_cSld.insert(0, copy.deepcopy(src_bg))
    # copy every shape from source spTree into new spTree
    new_spTree = new_slide.shapes._spTree
    src_spTree = src_slide.shapes._spTree
    for child in src_spTree:
        tag = child.tag
        if tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue  # keep the new tree's own group props
        new_spTree.append(copy.deepcopy(child))
    _remap_rels(src_slide, new_slide)
    return new_slide

# ---------------------------------------------------------------- text setting

def _set_runs(para, runs, default_font, default_size, default_color):
    """Replace a paragraph's runs with the given run specs.
    Each run spec is a dict: {text, font?, size?, color?, bold?, italic?}."""
    # clear existing runs
    for r in list(para.runs):
        r._r.getparent().remove(r._r)
    for spec in runs:
        r = para.add_run()
        r.text = spec["text"]
        f = r.font
        f.name = spec.get("font", default_font)
        sz = spec.get("size", default_size)
        if sz is not None:
            f.size = Pt(sz)
        f.bold = spec.get("bold", False)
        if spec.get("italic"):
            f.italic = True
        col = spec.get("color", default_color)
        if isinstance(col, str):
            col = COLMAP[col]
        if col is not None:
            f.color.rgb = col


def _clear_extra_paragraphs(tf):
    """Keep only the first paragraph element; remove the rest."""
    txBody = tf._txBody
    paras = txBody.findall(qn("a:p"))
    for p in paras[1:]:
        txBody.remove(p)


def set_block(shape, paragraphs, default_font, default_size, default_color):
    """Fill a text shape with a list of paragraph specs.
    Each paragraph spec is either:
      - a string (single default-styled run), or
      - a dict {runs:[...], align?, size?, space_after?} , or
      - None / "" for a blank spacer paragraph.
    """
    tf = shape.text_frame
    _clear_extra_paragraphs(tf)
    first = tf.paragraphs[0]
    # normalize to list
    for i, pspec in enumerate(paragraphs):
        para = first if i == 0 else tf.add_paragraph()
        if pspec is None or pspec == "":
            # blank spacer — clear runs, leave empty
            for r in list(para.runs):
                r._r.getparent().remove(r._r)
            if para.text:
                para.text = ""
            continue
        if isinstance(pspec, str):
            runs = [{"text": pspec}]
            align = None
        else:
            runs = pspec.get("runs", [{"text": pspec.get("text", "")}])
            align = pspec.get("align")
        psize = pspec.get("size") if isinstance(pspec, dict) else None
        _set_runs(para, runs, default_font,
                  psize if psize else default_size, default_color)
        if align == "center":
            from pptx.enum.text import PP_ALIGN
            para.alignment = PP_ALIGN.CENTER


def set_single(shape, text, **overrides):
    """Set a shape that holds one styled line, preserving its slot styling but
    replacing the text (and optional per-run overrides)."""
    tf = shape.text_frame
    _clear_extra_paragraphs(tf)
    para = tf.paragraphs[0]
    # capture existing first-run style as defaults
    runs = para.runs
    if runs:
        f0 = runs[0].font
        dfont = f0.name
        dsize = f0.size.pt if f0.size else None
        try:
            dcolor = f0.color.rgb if f0.color and f0.color.type == 1 else None
        except Exception:
            dcolor = None
        dbold = f0.bold
    else:
        dfont, dsize, dcolor, dbold = None, None, None, None
    spec = {"text": text,
            "font": overrides.get("font", dfont),
            "size": overrides.get("size", dsize),
            "color": overrides.get("color", dcolor),
            "bold": overrides.get("bold", dbold)}
    _set_runs(para, [spec], dfont, dsize, dcolor)

# ---------------------------------------------------------------- slot mapping
# Index of each editable shape within a cloned variant slide (clone preserves
# the template's shape order). Verified against the template inspection.
SLOTS = {
    1: {"page": 1, "eyebrow": 2, "title": 4, "subtitle": 5},
    2: {"page": 3, "eyebrow": 4, "title": 5, "body": 7},
    3: {"page": 1, "eyebrow": 2, "quote": 4, "attrib": 6},
    4: {"page": 1, "eyebrow": 3, "title": 5,
        "col1_head": 5 + 0, "col1": 6, "col2_head": 7, "col2": 8},
}
# NOTE: V4 title is shape 3; col1_head is shape 5 (corrected below in render).
SLOTS[4] = {"page": 1, "eyebrow": 2, "title": 3,
            "col1_head": 5, "col1": 6, "col2_head": 7, "col2": 8}

from pptx.enum.text import PP_ALIGN


def _para_to_spec(item, base_size, base_color, lead_color, lead_size):
    """Convert a body mini-language item into a set_block paragraph spec."""
    if item is None or item == "":
        return ""
    if isinstance(item, str):
        return {"runs": [{"text": item, "font": "Calibri",
                          "size": base_size, "color": base_color}]}
    # dict forms
    if "runs" in item:
        runs = []
        for r in item["runs"]:
            runs.append({"text": r["t"],
                         "font": r.get("font", "Calibri"),
                         "size": r.get("size", base_size),
                         "bold": r.get("b", False),
                         "color": r.get("c", base_color),
                         "italic": r.get("i", False)})
        out = {"runs": runs}
        if item.get("align") == "center":
            out["align"] = "center"
        return out
    if "lead" in item:
        lc = item.get("lead_c", lead_color)
        runs = [{"text": item["lead"], "font": "Georgia", "bold": True,
                 "size": item.get("lead_size", lead_size), "color": lc}]
        if item.get("text"):
            runs.append({"text": item["text"], "font": "Calibri",
                         "size": item.get("size", base_size),
                         "color": item.get("c", base_color)})
        return {"runs": runs}
    # plain text with optional color/size
    return {"runs": [{"text": item.get("text", ""), "font": "Calibri",
                      "size": item.get("size", base_size),
                      "bold": item.get("b", False),
                      "color": item.get("c", base_color)}]}


def fill_block(shape, items, base_size, base_color, lead_color, lead_size):
    specs = [_para_to_spec(it, base_size, base_color, lead_color, lead_size)
             for it in items]
    set_block(shape, specs, "Calibri", base_size, base_color)


def render_slide(prs, src_slides, spec, num):
    """src_slides: [V1,V2,V3,V4] template slides. spec: content dict. num: page #."""
    v = spec["v"]
    s = duplicate_slide(prs, src_slides[v - 1])
    shp = list(s.shapes)
    slot = SLOTS[v]
    set_single(shp[slot["page"]], f"{num:02d}")
    if "eyebrow" in spec:
        set_single(shp[slot["eyebrow"]], spec["eyebrow"])
    if v == 1:
        set_single(shp[slot["title"]], spec["title"])
        set_single(shp[slot["subtitle"]], spec.get("subtitle", ""))
    elif v == 2:
        set_single(shp[slot["title"]], spec["title"])
        bs = spec.get("body_size", 15)
        fill_block(shp[slot["body"]], spec.get("body", []),
                   bs, NAVY, TERRA, spec.get("lead_size", bs + 1))
    elif v == 3:
        set_single(shp[slot["quote"]], spec["quote"])
        set_single(shp[slot["attrib"]], spec.get("attrib", ""))
    elif v == 4:
        set_single(shp[slot["title"]], spec["title"])
        set_single(shp[slot["col1_head"]], spec.get("col1_head", ""))
        set_single(shp[slot["col2_head"]], spec.get("col2_head", ""),
                   color=GREEN)
        cs = spec.get("col_size", 13)
        fill_block(shp[slot["col1"]], spec.get("col1", []), cs, NAVY, TERRA, cs)
        fill_block(shp[slot["col2"]], spec.get("col2", []), cs, NAVY, GREEN, cs)
    return s


def build_deck(slides, out_path):
    prs = Presentation(TEMPLATE)
    src = list(prs.slides)              # [V1,V2,V3,V4]
    for i, spec in enumerate(slides, start=1):
        render_slide(prs, src, spec, i)
    # remove the 4 original template slides (now at the front)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    for el in ids[:4]:
        rId = el.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(el)
    prs.save(out_path)
    return len(slides)
