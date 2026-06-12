import deckbuild as db
from pptx import Presentation
from pptx.oxml.ns import qn

SLIDES = [
  # 1 — V1 divider / intro
  {"v":1, "eyebrow":"AUTOMATE THE ADMIN · GOING FURTHER",
   "title":"Claude Code — when chatting isn’t enough",
   "subtitle":"The same Claude you know, now working directly inside your files, folders and tools."},

  # 2 — V4 the difference: Chat vs Code
  {"v":4, "eyebrow":"TWO WAYS TO WORK WITH CLAUDE",
   "title":"Claude Chat vs. Claude Code",
   "col1_head":"Claude Chat",
   "col1":[
     {"lead":"It talks  ","text":"you type, it replies in the chat window."},
     {"lead":"You move the work  ","text":"copy text in, copy the answer back out by hand."},
     {"lead":"Best for  ","text":"thinking, drafting, summarising, answering questions."},
     "",
     {"text":"Great for words — you still do the doing.","c":"mgold"},
   ],
   "col2_head":"Claude Code",
   "col2":[
     {"lead":"It acts  ","text":"reads and edits your real files and runs the task itself."},
     {"lead":"It moves the work  ","text":"opens folders, edits many documents at once."},
     {"lead":"Best for  ","text":"repetitive, multi-file, “just change it” jobs."},
     "",
     {"text":"Great for doing — it works where your files live.","c":"mgold"},
   ]},

  # 3 — V2 how to use & integrate
  {"v":2, "eyebrow":"GETTING STARTED WITH CLAUDE CODE",
   "title":"How to use it — and where it fits",
   "body":[
     {"lead":"Point it at a folder  —  ","text":"open Claude Code where the files you want to work with already live."},
     {"lead":"Ask in plain English  —  ","text":"“rename these by date,” “pull the totals into one sheet,” “draft a report from these notes.”"},
     {"lead":"Watch and approve  —  ","text":"it shows each change before it makes it — nothing happens without your okay."},
     {"lead":"Use it where you work  —  ","text":"on your desktop (Mac or Windows), in the editor, in the terminal, or on the web."},
     {"lead":"Connect your tools  —  ","text":"link Google Drive, Gmail and more so it works with the data you already have."},
     {"text":"Start with one repetitive, file-heavy task — let it prove itself there first.","c":"mgold"},
   ], "body_size":14},
]

prs = Presentation(db.TEMPLATE)
src = list(prs.slides)
for i, spec in enumerate(SLIDES, start=1):
    s = db.render_slide(prs, src, spec, i)
    # blank the page number — these get inserted manually at unknown positions
    db.set_single(list(s.shapes)[db.SLOTS[spec["v"]]["page"]], "")
# drop the 4 template demo slides at the front
sldIdLst = prs.slides._sldIdLst
for el in list(sldIdLst)[:4]:
    prs.part.drop_rel(el.get(qn("r:id"))); sldIdLst.remove(el)
OUT="/Users/arjun/Desktop/IAF-Session1/IAF_claude_code_slides.pptx"
prs.save(OUT)
print("built", len(SLIDES), "slides ->", OUT)
