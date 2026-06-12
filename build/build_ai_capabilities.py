import deckbuild as db
from pptx import Presentation
from pptx.oxml.ns import qn

SLIDES = [
  # 1 — V3 framing: AI is the new electricity
  {"v":3, "eyebrow":"WHAT AI ACTUALLY IS",
   "quote":"A century ago, electricity went from a curiosity to the quiet power behind every industry. AI is on the same path — a general-purpose capability you plug into almost any kind of work.",
   "attrib":"“AI is the new electricity.” — Andrew Ng"},

  # 2 — V2 the hype (flows out of the Ng quote, into "what it actually is")
  {"v":2, "eyebrow":"FIRST, A REALITY CHECK",
   "title":"Also the most hyped word of the decade",
   "body":[
     {"lead":"Everything is “AI” now  —  ","text":"toothbrushes, spreadsheets, dating apps — the label sells, so it’s everywhere."},
     {"lead":"The promises run wild  —  ","text":"“replace your whole team,” “10× overnight,” “no skills needed.”"},
     {"lead":"Money pours in  —  ","text":"billions chase anything with “AI” in the pitch deck."},
     "",
     {"runs":[{"t":"Some of it is real. A lot of it is noise. ","b":True,"c":"terra"},{"t":"The skill is telling them apart — which is what today is about.","c":"navy"}]},
     {"text":"Don’t buy the magic wand — and don’t dismiss the real thing either.","c":"mgold"},
   ], "body_size":14},

  # 3 — V4 the two capabilities ("what it actually is")
  {"v":4, "eyebrow":"TWO KINDS OF CAPABILITY",
   "title":"Predictive AI vs. Generative AI",
   "col1_head":"Predictive",
   "col1":[
     {"lead":"What  ","text":"learns from past data to forecast and classify."},
     {"lead":"Asks  ","text":"“what’s likely?” — fraud, no-shows, a diagnosis."},
     {"lead":"In the wild  ","text":"predicting how proteins fold to design new medicines."},
     "",
     {"text":"Strongest with good data and a clear question.","c":"mgold"},
   ],
   "col2_head":"Generative",
   "col2":[
     {"lead":"What  ","text":"makes something new from a prompt — text, images, code."},
     {"lead":"Asks  ","text":"“make me one” — a draft, a summary, a picture."},
     {"lead":"In the wild  ","text":"Google’s “Nano Banana” turns a sentence into an image."},
     "",
     {"text":"You steer; it produces the first draft.","c":"mgold"},
   ], "col_size":12},

  # 4 — V1 sparse, image-ready: Isomorphic Labs (add website screenshot here)
  {"v":1, "eyebrow":"PREDICTIVE AI · REAL-WORLD IMPACT",
   "title":"Predicting the building blocks of life",
   "subtitle":"Isomorphic Labs uses AI to predict how proteins fold — turning years of lab work into a problem you can solve on a computer, and speeding the search for new medicines."},

  # 5 — V1 sparse, image-ready: AI does the impossible (Vesuvius scrolls — swap if desired)
  {"v":1, "eyebrow":"WHEN AI DOES THE IMPOSSIBLE",
   "title":"Reading a scroll no human could open",
   "subtitle":"Buried by Vesuvius in 79 AD and too charred to unroll, the Herculaneum scrolls sat unreadable for 2,000 years. In 2024, AI virtually unrolled one and recovered its Greek text — without ever touching it."},

  # 6 — V1 sparse, image-ready: the grift (add the ₹9 course screenshot + one more)
  {"v":1, "eyebrow":"AND ALSO… THIS",
   "title":"Where there’s hype, there’s a hustle",
   "subtitle":"“Become an AI millionaire — only ₹9.” Hype grows its own economy of shortcuts and gurus. Real skill isn’t for sale at a discount — it’s built by doing, which is exactly why you’re here."},
]

prs = Presentation(db.TEMPLATE)
src = list(prs.slides)
for i, spec in enumerate(SLIDES, start=1):
    s = db.render_slide(prs, src, spec, i)
    # blank the page number — these get inserted manually at unknown positions
    db.set_single(list(s.shapes)[db.SLOTS[spec["v"]]["page"]], "")
# drop the 4 template demo slides at the front
sldIdLst = prs.slides._sldIdLst
for el in list(sldIdLst)[:4]:
    prs.part.drop_rel(el.get(qn("r:id"))); sldIdLst.remove(el)
OUT="/Users/arjun/Desktop/IAF-Session1/IAF_ai_capabilities_slides.pptx"
prs.save(OUT)
print("built", len(SLIDES), "slides ->", OUT)
