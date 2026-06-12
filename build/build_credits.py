import deckbuild as db
from pptx import Presentation
from pptx.oxml.ns import qn

SLIDES = [
  # 1 — V1 divider
  {"v":1, "eyebrow":"WORKING SUSTAINABLY",
   "title":"Make your credits last",
   "subtitle":"Match the model to the task — and spend tokens like the budget they are."},

  # 2 — V4 big vs small model
  {"v":4, "eyebrow":"RIGHT-SIZE THE MODEL",
   "title":"Big model, or small model?",
   "col1_head":"Reach for the powerful one",
   "col1":[
     {"lead":"Use it for  ","text":"hard reasoning, ambiguous problems, gnarly code, the tricky first draft."},
     "",
     {"lead":"Such as  ","text":"Claude Opus · ChatGPT’s deep-reasoning mode."},
     "",
     {"text":"Powerful models cost more per message — save them for when they earn it.","c":"mgold"},
   ],
   "col2_head":"Default to the fast one",
   "col2":[
     {"lead":"Use it for  ","text":"everyday drafting, summaries, reformatting, repetitive Q&A, quick edits."},
     "",
     {"lead":"Such as  ","text":"Claude Sonnet or Haiku · ChatGPT’s fast mode."},
     "",
     {"text":"Most of your day-to-day work never needs the expensive model.","c":"mgold"},
   ]},

  # 3 — V2 habits
  {"v":2, "eyebrow":"HABITS THAT SAVE TOKENS",
   "title":"Spend less, get more",
   "body":[
     {"lead":"Start a fresh chat for a new task  —  ","text":"long histories re-send every earlier message, and you pay for them each time."},
     {"lead":"Be specific, not verbose  —  ","text":"a clear, short prompt beats a long, rambling one."},
     {"lead":"Don’t paste the whole document  —  ","text":"give the one section the model actually needs."},
     {"lead":"Reuse your saved prompts  —  ","text":"don’t re-explain the same setup from scratch."},
     {"lead":"Turn off tools you’re not using  —  ","text":"web search and extras add tokens you may not need."},
     {"text":"Every token you don’t spend is budget for another month of building.","c":"mgold"},
   ], "body_size":15},

  # 4 — V3 rule of thumb
  {"v":3, "eyebrow":"THE RULE OF THUMB",
   "quote":"Start with the smallest model that could do the job. Escalate only when it visibly struggles — most everyday work never needs the expensive one.",
   "attrib":"Right-sizing is the cheapest habit you’ll build"},
]

prs = Presentation(db.TEMPLATE)
src = list(prs.slides)
for i, spec in enumerate(SLIDES, start=1):
    s = db.render_slide(prs, src, spec, i)
    # blank the page number — these get inserted manually at unknown positions
    db.set_single(list(s.shapes)[db.SLOTS[spec["v"]]["page"]], "")
# drop the 4 template demo slides at the front
sldIdLst = prs.slides._sldIdLst
for el in list(sldIdLst)[:4]:
    prs.part.drop_rel(el.get(qn("r:id"))); sldIdLst.remove(el)
OUT="/Users/arjun/Desktop/IAF-Session1/IAF_credits_tokens_slides.pptx"
prs.save(OUT)
print("built", len(SLIDES), "slides ->", OUT)
